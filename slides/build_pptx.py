"""Build the ICMIC 2026 PowerPoint deck for CONFROUTE.

One idea per slide, image-led. Charts come from make_figures.py (PNG at 200 dpi),
photographs from fetch_images.py (Wikimedia Commons, credited on the last slide),
satellite chips from the ShipsNet dataset.

    python slides/make_figures.py && python slides/build_pptx.py
"""

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from PIL import Image
from lxml import etree

HERE = Path(__file__).resolve().parent
FIG = HERE / "figures"
IMG = HERE / "images"
ARCH = HERE.parent / "paper" / "figures" / "architecture_full.png"
OUT = HERE / "CONFROUTE_ICMIC2026.pptx"

W, H = Inches(13.333), Inches(7.5)

# palette, identical to the figures
NAVY = RGBColor(0x10, 0x46, 0x80)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
AQUA = RGBColor(0x1B, 0xAF, 0x7A)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
RULE = RGBColor(0xE1, 0xE0, 0xD9)
WASH = RGBColor(0xF4, 0xF7, 0xFB)
WASH_O = RGBColor(0xFD, 0xF0, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"
MARGIN = Inches(0.72)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------------------- #
# primitives
# --------------------------------------------------------------------------- #
def slide(notes=None):
    s = prs.slides.add_slide(BLANK)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def _alpha(color_elm, alpha):
    """python-pptx has no transparency API; splice an <a:alpha> into the fill."""
    srgb = color_elm.find(qn("a:srgbClr"))
    if srgb is None:
        return
    a = etree.SubElement(srgb, qn("a:alpha"))
    a.set("val", str(int(alpha * 100000)))


def rect(s, x, y, w, h, color, alpha=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if alpha < 1.0:
        _alpha(shp.fill._xPr.find(qn("a:solidFill")), alpha)
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line=1.15, space_after=0):
    """runs: list of (string, size_pt, bold, color) or a list of such lists (paragraphs)."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(space_after)
        for content, size, bold, color in para:
            r = p.add_run()
            r.text = content
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = FONT
    return box


def pic_cover(s, path, x, y, w, h):
    """Fill the box, preserving aspect ratio by centre-cropping (CSS object-fit: cover)."""
    iw, ih = Image.open(path).size
    box_ar, img_ar = w / h, iw / ih
    pic = s.shapes.add_picture(str(path), x, y, w, h)
    if img_ar > box_ar:                       # image too wide -> crop sides
        keep = box_ar / img_ar
        pic.crop_left = pic.crop_right = (1 - keep) / 2
    else:                                     # image too tall -> crop top/bottom
        keep = img_ar / box_ar
        pic.crop_top = pic.crop_bottom = (1 - keep) / 2
    return pic


def pic_fit(s, path, x, y, w, h):
    """Contain the whole image inside the box, centred."""
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    nw, nh = int(iw * scale), int(ih * scale)
    return s.shapes.add_picture(str(path), x + int((w - nw) / 2), y + int((h - nh) / 2), nw, nh)


def footer(s, n):
    text(s, MARGIN, H - Inches(0.52), Inches(8), Inches(0.3),
         [("CONFROUTE  |  ICMIC 2026, NTUST Taipei", 10, False, MUTED)])
    text(s, W - MARGIN - Inches(1), H - Inches(0.52), Inches(1), Inches(0.3),
         [(str(n), 10, False, MUTED)], align=PP_ALIGN.RIGHT)


# --------------------------------------------------------------------------- #
# slide templates
# --------------------------------------------------------------------------- #
def photo_slide(headline, sub=None, image=None, notes=None, accent=None, side="bottom"):
    """Full-bleed photograph with a scrim and one short headline."""
    s = slide(notes)
    pic_cover(s, IMG / image, 0, 0, W, H)
    if side == "bottom":
        rect(s, 0, Inches(3.9), W, H - Inches(3.9), RGBColor(0x05, 0x14, 0x24), 0.72)
        y = Inches(4.55)
    else:
        rect(s, 0, 0, W, H, RGBColor(0x05, 0x14, 0x24), 0.55)
        y = Inches(2.7)
    text(s, MARGIN, y, W - 2 * MARGIN, Inches(1.6),
         [(headline, 40, True, WHITE)], line=1.1)
    if sub:
        gap = Inches(1.35) if side == "bottom" else Inches(1.75)
        text(s, MARGIN, y + gap, W - 2 * MARGIN, Inches(1.0),
             [(sub, 20, False, RGBColor(0xD8, 0xE4, 0xF2))], line=1.25)
    if accent:
        rect(s, MARGIN, y - Inches(0.34), Inches(1.1), Pt(5), accent or BLUE)
    return s


def statement_slide(headline, sub=None, notes=None, color=NAVY):
    """White slide, one sentence, nothing else."""
    s = slide(notes)
    text(s, Inches(1.4), Inches(2.5), W - Inches(2.8), Inches(1.9),
         [(headline, 44, True, color)], align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.BOTTOM, line=1.12)
    if sub:
        text(s, Inches(2.1), Inches(4.6), W - Inches(4.2), Inches(1.4),
             [(sub, 20, False, INK2)], align=PP_ALIGN.CENTER, line=1.35)
    return s


def figure_slide(title, image, n, sub=None, note=None, notes=None, wide=0.86):
    """A chart or diagram with the takeaway as the title."""
    s = slide(notes)
    text(s, MARGIN, Inches(0.5), W - 2 * MARGIN, Inches(0.8),
         [(title, 30, True, NAVY)], line=1.1)
    top = Inches(1.32)
    if sub:
        text(s, MARGIN, Inches(1.22), W - 2 * MARGIN, Inches(0.45),
             [(sub, 16, False, INK2)])
        top = Inches(1.78)
    bottom = H - Inches(1.35) if note else H - Inches(0.75)
    pic_fit(s, image, MARGIN, top, W - 2 * MARGIN, bottom - top)
    if note:
        band = rect(s, MARGIN, H - Inches(1.24), W - 2 * MARGIN, Inches(0.72), WASH_O)
        text(s, MARGIN + Inches(0.28), H - Inches(1.13), W - 2 * MARGIN - Inches(0.56),
             Inches(0.5), note, anchor=MSO_ANCHOR.MIDDLE, line=1.2)
        _ = band
    footer(s, n)
    return s


def title_only(title, n, sub=None, notes=None):
    s = slide(notes)
    text(s, MARGIN, Inches(0.5), W - 2 * MARGIN, Inches(0.8),
         [(title, 30, True, NAVY)], line=1.1)
    if sub:
        text(s, MARGIN, Inches(1.22), W - 2 * MARGIN, Inches(0.45),
             [(sub, 16, False, INK2)])
    footer(s, n)
    return s


# --------------------------------------------------------------------------- #
# the deck
# --------------------------------------------------------------------------- #
n = 0


def num():
    global n
    n += 1
    return n


# ---- 1. title -------------------------------------------------------------
s = slide("Name, affiliation, one line: this is a paper about when an edge AI "
          "system on a ship should refuse to make a decision. 20 seconds, do not linger.")
pic_cover(s, IMG / "hero.jpg", 0, 0, W, H)
rect(s, 0, 0, W, H, RGBColor(0x04, 0x11, 0x20), 0.62)
text(s, MARGIN, Inches(1.05), W - 2 * MARGIN, Inches(0.4),
     [("ICMIC 2026  ·  NTUST, Taipei  ·  19–21 August 2026", 14, False,
       RGBColor(0x9E, 0xC5, 0xF4))])
text(s, MARGIN, Inches(1.95), W - 2 * MARGIN, Inches(1.1),
     [("CONFROUTE", 66, True, WHITE)], line=1.0)
text(s, MARGIN, Inches(3.0), W - Inches(4.0), Inches(1.3),
     [("Confidence-Aware Routing for Resilient Maritime Edge AI", 28, False,
       RGBColor(0xCD, 0xE2, 0xFB))], line=1.2)
rect(s, MARGIN, Inches(4.45), Inches(1.6), Pt(5), BLUE)
text(s, MARGIN, Inches(4.95), W - 2 * MARGIN, Inches(0.9),
     [[("Anthony Uchenna Eneh", 19, True, WHITE),
       ("   ·   Love Allen Chijioke Ahakonye   ·   Jae Min Lee   ·   Dong-Seong Kim",
        19, False, RGBColor(0xD8, 0xE4, 0xF2))]], line=1.3)
text(s, MARGIN, Inches(5.65), W - 2 * MARGIN, Inches(1.0),
     [("Kumoh National Institute of Technology, Gumi, South Korea   ·   NSLab Co. Ltd.",
       14, False, RGBColor(0x9E, 0xC5, 0xF4))], line=1.35)

# ---- 2-4. the operating environment, one constraint each ------------------
photo_slide("Ships now run AI on board.",
            "Perception, monitoring, video triage, navigation support.",
            image="sensor.jpg", accent=BLUE,
            notes="Set the scene fast. Edge AI is already deployed at sea. "
                  "Do not list applications, the slide does it. ~25 s.")
footer(prs.slides[-1], num())

photo_slide("The link is not always there.",
            "Satellite and ship-to-ship connectivity is intermittent by nature.",
            image="satcom.jpg", accent=BLUE,
            notes="Constraint one. Offloading is not a guaranteed escape hatch. ~25 s.")
footer(prs.slides[-1], num())

photo_slide("And the sea changes what the camera sees.",
            "Sea state, illumination, haze, occlusion: the input drifts away from training.",
            image="fog.jpg", accent=ORANGE,
            notes="Constraint two, and the one this paper is really about. "
                  "Point at the haze. ~25 s.")
footer(prs.slides[-1], num())

# ---- 5. the stakes --------------------------------------------------------
photo_slide("A wrong answer costs more than a slow one.",
            "That inversion is what separates this from ordinary edge offloading.",
            image="storm.jpg", accent=ORANGE, side="full",
            notes="The key inversion. Latency-optimal is not safety-optimal. "
                  "Say it plainly and pause. ~25 s.")
footer(prs.slides[-1], num())

# ---- 6. the question ------------------------------------------------------
statement_slide("Act now, ask for help, or refuse to answer?",
                "Every task, on every node, under uncertainty.",
                notes="The question the paper answers. Let it sit for a beat. ~15 s.")
footer(prs.slides[-1], num())

# ---- 7. the assumption ----------------------------------------------------
statement_slide("Everyone assumes the model\nknows when it is unsure.",
                "Confidence-threshold offloading is built entirely on that assumption.",
                notes="Prior maritime MEC work optimises latency against energy and "
                      "treats confidence as a trustworthy gate. We tested it. ~25 s.")
footer(prs.slides[-1], num())

# ---- 8-9. the measurement, on real chips ---------------------------------
figure_slide("Clear water: the classifier finds every ship.",
             FIG / "fig_chips_nominal.png", num(),
             sub="Three held-out ShipsNet chips, 80 × 80 PlanetScope imagery.",
             notes="Real chips from the dataset. All three correct. "
                   "Note confidence is only 0.51 to 0.67: honestly unsure. ~30 s.")

figure_slide("Add haze, and it stops seeing them.",
             FIG / "fig_chips.png", num(),
             note=[("Every chip still contains a ship. The model is now ", 15, False, INK),
                   ("wrong on all three, and far more confident", 15, True, ORANGE),
                   (".", 15, False, INK)],
             notes="THE slide. Same three vessels after the adverse transform. "
                   "0.67 correct becomes 0.96 WRONG. Slow down here. ~45 s.")

# ---- 10. the aggregate ----------------------------------------------------
figure_slide("It is not three unlucky chips.",
             FIG / "fig_calibration.png", num(),
             sub="500 held-out images, nominal versus adverse.",
             note=[("Accuracy falls 13.8 points. Mean confidence ", 15, False, INK),
                   ("rises", 15, True, ORANGE),
                   (" 3.2 points.", 15, False, INK)],
             notes="The population-level version of the previous slide. "
                   "The lines cross. That crossing is the paper. ~40 s.")

# ---- 11. the consequence --------------------------------------------------
figure_slide("So a confidence gate cannot see the errors.",
             FIG / "fig_error_confidence.png", num(),
             note=[("Under adverse conditions ", 15, False, INK),
                   ("90% of errors sit above the 0.65 threshold", 15, True, ORANGE),
                   (", invisible to it.", 15, False, INK)],
             notes="Errors only, plotted against the fallback threshold. "
                   "39% escape under nominal, 90% under adverse. ~40 s.")

statement_slide("A confidence gate fails\nexactly when you need it.",
                notes="One line, then move to the method. ~12 s.", color=ORANGE)
footer(prs.slides[-1], num())

# ---- 13. the method -------------------------------------------------------
statement_slide("CONFROUTE",
                "Make refusal a first-class routing action, alongside compute and offload.",
                notes="Section break into the method. ~15 s.")
footer(prs.slides[-1], num())

figure_slide("One layer, between the model and the actuator.",
             ARCH, num(),
             sub="Consumes uncertainty without touching the model. Logs every decision for audit.",
             notes="Trace the three branches with the pointer. The router is post-scoring, "
                   "pre-actuation, so it drops into an existing stack. ~40 s.")

# ---- 15. three actions ----------------------------------------------------
s = slide("Local when confident and the queue allows. Offload when the link is good "
          "and the queue is loaded. Fallback when uncertain AND cut off. ~35 s.")
text(s, MARGIN, Inches(0.5), W - 2 * MARGIN, Inches(0.8),
     [("Three actions, one of them new.", 30, True, NAVY)], line=1.1)
cards = [
    ("LOCAL", BLUE, "Run the model here.", "Confident enough to act alone."),
    ("OFFLOAD", AQUA, "Ask a peer or relay.", "Link is usable, queue is loaded."),
    ("FALLBACK", ORANGE, "Refuse, and act safely.", "Uncertain and cut off."),
]
cw, gap = Inches(3.75), Inches(0.42)
x0 = (W - (3 * cw + 2 * gap)) / 2
for i, (name, colour, line1, line2) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rect(s, x, Inches(1.85), cw, Inches(3.5), WASH if i < 2 else WASH_O)
    rect(s, x, Inches(1.85), cw, Pt(7), colour)
    text(s, x + Inches(0.4), Inches(2.35), cw - Inches(0.8), Inches(0.6),
         [(name, 26, True, colour)])
    text(s, x + Inches(0.4), Inches(3.15), cw - Inches(0.8), Inches(1.0),
         [(line1, 19, True, INK)], line=1.2)
    text(s, x + Inches(0.4), Inches(3.95), cw - Inches(0.8), Inches(1.2),
         [(line2, 16, False, INK2)], line=1.25)
text(s, MARGIN, Inches(5.75), W - 2 * MARGIN, Inches(0.6),
     [[("Scores decide, ", 17, False, INK2),
       ("guardrails override", 17, True, NAVY),
       (": no offload without a peer, no local execution on a safety-critical task "
        "below the confidence floor.", 17, False, INK2)]], line=1.3)
footer(s, num())

# ---- 16. fallback is concrete --------------------------------------------
s = title_only("Fallback is an action, not an error code.", num(),
               sub="What the conservative branch means, per maritime application.",
               notes="Kills the objection that fallback means crash or do nothing. "
                     "Read the collision-avoidance row aloud, then move on. ~35 s.")
rows = [
    ("Collision avoidance", "Reduce speed, sound horn, hold heading, alert the operator"),
    ("Maritime surveillance", "Flag the chip for review, defer the alert, request another pass"),
    ("Engine anomaly", "Report “inspection needed”, fall back to a limit-based alert"),
    ("Waypoint decision", "Hold position, revert to last safe heading, request human review"),
]
y = Inches(2.05)
for i, (app, act) in enumerate(rows):
    if i % 2 == 0:
        rect(s, MARGIN, y - Inches(0.12), W - 2 * MARGIN, Inches(0.95), WASH)
    text(s, MARGIN + Inches(0.3), y, Inches(3.6), Inches(0.7),
         [(app, 19, True, NAVY)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, MARGIN + Inches(4.1), y, W - 2 * MARGIN - Inches(4.5), Inches(0.7),
         [(act, 18, False, INK)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.06)
text(s, MARGIN, Inches(6.42), W - 2 * MARGIN, Inches(0.5),
     [("Conservative, auditable, and aligned with COLREGS/IMO practice: "
       "the ship does something safe and says why.", 15, False, INK2)])

# ---- 17. when fallback fires ---------------------------------------------
figure_slide("Refuse only when both signals fail.",
             FIG / "fig_boundary.png", num(),
             sub="Confidence below 0.65 and bandwidth below 1.8 Mbps.",
             note=[("If the link is healthy, the router would rather ", 15, False, INK),
                   ("ask a better-resourced peer than refuse", 15, True, NAVY),
                   (". That keeps fallback rare.", 15, False, INK)],
             notes="Requiring both conditions is what stops the router hiding behind "
                   "fallback. ~30 s.")

# ---- 18. setup ------------------------------------------------------------
s = slide("ShipsNet, ResNet-18, three link regimes, four baselines. Say the unsafe-rate "
          "definition out loud: routing-layer risk, not physical incidents. ~35 s.")
text(s, MARGIN, Inches(0.5), W - 2 * MARGIN, Inches(0.8),
     [("37,500 routing decisions.", 30, True, NAVY)], line=1.1)
text(s, MARGIN, Inches(1.22), W - 2 * MARGIN, Inches(0.45),
     [("2,500 tasks × 5 policies × 3 link regimes.", 16, False, INK2)])
stats = [
    ("4,000", "PlanetScope chips\n80 × 80, ShipsNet"),
    ("ResNet-18", "ImageNet-initialised\n84.8% validation"),
    ("3", "link regimes\nstable / intermittent / degraded"),
    ("4", "baselines\nlocal, offload, confidence, load"),
]
cw, gap = Inches(2.85), Inches(0.32)
x0 = (W - (4 * cw + 3 * gap)) / 2
for i, (big, small) in enumerate(stats):
    x = x0 + i * (cw + gap)
    rect(s, x, Inches(2.1), cw, Inches(2.6), WASH)
    text(s, x, Inches(2.5), cw, Inches(0.9), [(big, 40, True, BLUE)],
         align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.25), Inches(3.5), cw - Inches(0.5), Inches(1.1),
         [(small, 15, False, INK2)], align=PP_ALIGN.CENTER, line=1.3)
text(s, MARGIN, Inches(5.25), W - 2 * MARGIN, Inches(1.0),
     [[("Unsafe rate", 17, True, NAVY),
       (" = routing-layer risk: safety-sensitive local errors, low-confidence local "
        "decisions, failed safety-sensitive offloads.\nIt is not a count of physical "
        "incidents.", 17, False, INK2)]], line=1.35)
footer(s, num())

# ---- 19-23. results -------------------------------------------------------
figure_slide("Every baseline converges at about 7%.",
             FIG / "fig_unsafe_by_scenario.png", num(),
             note=[("When the link goes, so does their escape hatch. ", 15, False, INK),
                   ("CONFROUTE holds at 4.8%", 15, True, BLUE),
                   (".", 15, False, INK)],
             notes="The convergence is the story. Stable links flatter everything that "
                   "offloads; degraded links expose the dependency. ~45 s.")

figure_slide("Degraded links: safer and faster.",
             FIG / "fig_tradeoff.png", num(),
             note=[("Unsafe 6.8% → 4.8% and latency 98.2 → 58.5 ms against always-offload. ",
                    15, False, INK),
                   ("A failed offload costs both.", 15, True, BLUE)],
             notes="Normally safety costs latency. Not here, because attempting an "
                   "offload that fails wastes time AND leaves the task unhandled. ~40 s.")

s = title_only("Where we do not win: intermittent links.", num(),
               sub="The regime in which offloading still mostly works.",
               notes="Deliberate honesty, and it buys credibility. Load-only is safer "
                     "here. Say why, then land the closing line. ~40 s.")
tbl_rows = [
    ("Always local", "49.4", "5.8%", "0.0%", False),
    ("Always offload", "71.6", "4.3%", "0.0%", False),
    ("Confidence threshold", "59.6", "4.9%", "0.0%", False),
    ("Load only", "67.0", "3.8%", "0.0%", False),
    ("CONFROUTE (ours)", "61.6", "4.4%", "2.6%", True),
]
hx = [MARGIN + Inches(0.3), MARGIN + Inches(4.6), MARGIN + Inches(6.7), MARGIN + Inches(8.9)]
for label, x in zip(["Policy", "Latency (ms)", "Unsafe rate", "Fallback"], hx):
    text(s, x, Inches(2.0), Inches(4.2), Inches(0.4), [(label, 16, True, NAVY)])
rect(s, MARGIN + Inches(0.3), Inches(2.45), W - 2 * MARGIN - Inches(0.6), Pt(1.5), RULE)
y = Inches(2.62)
for label, lat, uns, fb, ours in tbl_rows:
    if ours:
        rect(s, MARGIN, y - Inches(0.1), W - 2 * MARGIN, Inches(0.62), WASH)
    c = BLUE if ours else INK
    for val, x in zip([label, lat, uns, fb], hx):
        text(s, x, y, Inches(4.2), Inches(0.5), [(val, 19, ours, c)])
    y += Inches(0.66)
rect(s, MARGIN, Inches(6.05), W - 2 * MARGIN, Inches(0.92), WASH_O)
text(s, MARGIN + Inches(0.3), Inches(6.2), W - 2 * MARGIN - Inches(0.6), Inches(0.7),
     [[("Load-only routing is safer here", 16, True, INK),
       (" (3.8% vs 4.4%), at 5.4 ms more latency. Our case is the degraded regime, "
        "which is also the regime that actually threatens the vessel.", 16, False, INK)]],
     line=1.3)

figure_slide("The router shifts with the link.",
             FIG / "fig_actions.png", num(),
             note=[("Offload collapses 89% → 11%. Fallback stays a ", 15, False, INK),
                   ("bounded minority", 15, True, NAVY),
                   (": the vessel keeps operating.", 15, False, INK)],
             notes="Shows adaptation, and that fallback never runs away. ~30 s.")

s = title_only("Both parts are load-bearing.", num(),
               sub="Unsafe outcome rate under degraded links.",
               notes="Removing fallback hurts more than removing link awareness. "
                     "Having somewhere to put an undecidable task is the bigger effect. ~30 s.")
abl = [("Full CONFROUTE", "4.8%", "", True),
       ("Without the fallback action", "6.2%", "p = 0.012", False),
       ("Without link-state awareness", "5.9%", "p = 0.037", False)]
y = Inches(2.5)
for label, val, sig, ours in abl:
    if ours:
        rect(s, MARGIN, y - Inches(0.16), W - 2 * MARGIN, Inches(0.86), WASH)
    text(s, MARGIN + Inches(0.35), y, Inches(6.5), Inches(0.6),
         [(label, 22, ours, BLUE if ours else INK)])
    text(s, MARGIN + Inches(7.2), y, Inches(2.0), Inches(0.6),
         [(val, 26, True, BLUE if ours else INK)])
    text(s, MARGIN + Inches(9.4), y, Inches(2.5), Inches(0.6),
         [(sig, 19, False, INK2)])
    y += Inches(1.05)

# ---- 24. the point --------------------------------------------------------
statement_slide("The routing layer was\nnever the bottleneck.",
                "We improved it, and unsafe outcomes still occur at high confidence.",
                notes="Slow down. Let it sit before speaking. This is what you want "
                      "the room to remember. ~25 s.")
footer(prs.slides[-1], num())

# ---- 25. limitations ------------------------------------------------------
s = title_only("What we are not claiming.", num(),
               notes="State the limitations without apologising, then the forward path "
                     "with calibration first. ~40 s.")
cols = [
    # one line each: wrapped bullets lose their hanging indent
    ("Limitations", ORANGE, [
        "Degradation is synthetic.",
        "One task family: satellite chips.",
        "Outcomes from a cost model.",
        "Thresholds fixed and trace-tuned.",
    ]),
    ("Where this goes", NAVY, [
        "Calibration first: conformal, OOD.",
        "Adaptive thresholds via bandits.",
        "Peer-aware forwarding.",
        "Real traces, hardware-in-the-loop.",
    ]),
]
cw = Inches(5.7)
for i, (head, colour, items) in enumerate(cols):
    x = MARGIN + i * (cw + Inches(0.7))
    text(s, x, Inches(1.75), cw, Inches(0.5), [(head, 22, True, colour)])
    rect(s, x, Inches(2.28), Inches(1.0), Pt(4), colour)
    y = Inches(2.65)
    for it in items:
        text(s, x, y, cw, Inches(1.0), [("•  " + it, 17, False, INK)], line=1.3)
        y += Inches(1.02)

# ---- 26. takeaways --------------------------------------------------------
s = title_only("Three things to take away.", num(),
               notes="Read the three headers only, then stop and invite questions. "
                     "Do not add a new thought here. ~40 s.")
takeaways = [
    ("Fallback deserves to be a routing action.",
     "Under degraded links it cuts unsafe outcomes from 6.8% to 4.8% and saves 40 ms."),
    ("Confidence thresholds break exactly when needed.",
     "Accuracy falls 13.8 points while confidence rises, and 90% of errors clear the gate."),
    ("Calibration is the reliability bottleneck.",
     "Ahead of routing or scheduling. Safe autonomy at sea depends on knowing "
     "what the model does not know."),
]
y = Inches(1.9)
for i, (head, body) in enumerate(takeaways):
    text(s, MARGIN, y, Inches(0.7), Inches(0.7), [(str(i + 1), 34, True, BLUE)])
    text(s, MARGIN + Inches(0.85), y + Inches(0.03), W - 2 * MARGIN - Inches(0.9),
         Inches(0.5), [(head, 23, True, INK)], line=1.15)
    text(s, MARGIN + Inches(0.85), y + Inches(0.62), W - 2 * MARGIN - Inches(0.9),
         Inches(0.8), [(body, 17, False, INK2)], line=1.3)
    y += Inches(1.55)

# ---- 27. thank you --------------------------------------------------------
s = slide("Thank you, and take questions. Backup slides follow.")
pic_cover(s, IMG / "hero.jpg", 0, 0, W, H)
rect(s, 0, 0, W, H, RGBColor(0x04, 0x11, 0x20), 0.66)
text(s, MARGIN, Inches(2.6), W - 2 * MARGIN, Inches(1.2),
     [("Thank you. Questions?", 48, True, WHITE)], line=1.1)
rect(s, MARGIN, Inches(3.95), Inches(1.6), Pt(5), BLUE)
text(s, MARGIN, Inches(4.4), W - 2 * MARGIN, Inches(1.6),
     [[("Anthony Uchenna Eneh\n", 20, True, WHITE),
       ("anthony@kumoh.ac.kr\nKumoh National Institute of Technology, Gumi, South Korea",
        18, False, RGBColor(0xD8, 0xE4, 0xF2))]], line=1.4)
num()

# ---- backup ---------------------------------------------------------------
s = slide()
text(s, Inches(1.4), Inches(3.2), W - Inches(2.8), Inches(1.0),
     [("Backup", 40, True, MUTED)], align=PP_ALIGN.CENTER)

s = title_only("Backup: full results", num())
head = ["Policy", "Regime", "Latency (ms)", "Unsafe", "Energy (J)"]
hx = [MARGIN + Inches(0.2), MARGIN + Inches(3.5), MARGIN + Inches(5.6),
      MARGIN + Inches(8.0), MARGIN + Inches(9.8)]
for label, x in zip(head, hx):
    text(s, x, Inches(1.55), Inches(3.4), Inches(0.35), [(label, 13, True, NAVY)])
rect(s, MARGIN + Inches(0.2), Inches(1.9), W - 2 * MARGIN - Inches(0.4), Pt(1.2), RULE)
full = [
    ("Always local", "Stable", "49.3", "6.7%", "0.50"), ("Always offload", "Stable", "33.4", "1.7%", "0.60"),
    ("Confidence threshold", "Stable", "46.5", "3.4%", "0.53"), ("Load only", "Stable", "48.2", "6.4%", "0.51"),
    ("CONFROUTE", "Stable", "33.5", "1.5%", "0.60"),
    ("Always local", "Intermittent", "49.4", "5.8%", "0.50"), ("Always offload", "Intermittent", "71.6", "4.3%", "0.59"),
    ("Confidence threshold", "Intermittent", "59.6", "4.9%", "0.53"), ("Load only", "Intermittent", "67.0", "3.8%", "0.58"),
    ("CONFROUTE", "Intermittent", "61.6", "4.4%", "0.56"),
    ("Always local", "Degraded", "49.3", "7.3%", "0.50"), ("Always offload", "Degraded", "98.2", "6.8%", "0.56"),
    ("Confidence threshold", "Degraded", "78.9", "7.0%", "0.55"), ("Load only", "Degraded", "91.8", "7.2%", "0.56"),
    ("CONFROUTE", "Degraded", "58.5", "4.8%", "0.50"),
]
y = Inches(2.02)
for row in full:
    ours = row[0] == "CONFROUTE"
    if ours:
        rect(s, MARGIN, y - Inches(0.04), W - 2 * MARGIN, Inches(0.36), WASH)
    for val, x in zip(row, hx):
        text(s, x, y, Inches(3.4), Inches(0.3), [(val, 13, ours, BLUE if ours else INK)])
    y += Inches(0.345)

s = title_only("Backup: confidence traces", num())
text(s, MARGIN, Inches(1.9), W - 2 * MARGIN, Inches(2.6),
     [[("Nominal", 22, True, INK), ("      500 samples      88.8% accuracy      "
       "mean confidence 0.810      std 0.140\n", 20, False, INK2),
       ("Adverse", 22, True, INK), ("      500 samples      75.0% accuracy      "
       "mean confidence 0.842      std 0.085\n\n", 20, False, INK2),
       ("The adverse transform is deterministic: reduced contrast and brightness, "
        "Gaussian blur, haze blending, sensor-like noise.\nNote the narrowing standard "
        "deviation: the model is not just wrong, it is uniformly confident.\n\n",
        17, False, INK2),
       ("Mean confidence on incorrect predictions: 0.651 nominal, 0.782 adverse.\n"
        "Errors above the 0.65 threshold: 22 of 56 nominal, 113 of 125 adverse.",
        17, True, NAVY)]], line=1.5)

# ---- credits --------------------------------------------------------------
# generated from the fetch manifest so attribution cannot drift from the files
s = title_only("Image credits", num(),
               sub="Photographs from Wikimedia Commons, reused under the licence shown.")
SLOT_DESC = {
    "hero": "Aerial container ship (title, closing)",
    "fog": "Vessel in haze",
    "storm": "Bulk carrier in heavy seas",
    "satcom": "Shipborne radomes",
    "sensor": "Frigate radar masts",
}
def tidy_author(raw):
    """Commons artist fields for derivative works carry filename noise."""
    a = " ".join((raw or "see source").split())
    a = re.sub(r"^File:.*?\.(jpg|jpeg|png|tif+)\s*:\s*", "", a, flags=re.I)
    a = re.sub(r"\s*derivative work\s*:\s*", " / ", a, flags=re.I)
    a = re.sub(r"^(User|user):\s*", "", a)
    return a[:60].strip()


y = Inches(2.0)
for c in json.loads((IMG / "CREDITS.json").read_text(encoding="utf-8")):
    author = tidy_author(c["artist"])
    text(s, MARGIN, y, Inches(4.6), Inches(0.4),
         [(SLOT_DESC.get(c["slot"], c["slot"]), 16, True, INK)])
    text(s, MARGIN + Inches(4.9), y, W - 2 * MARGIN - Inches(4.9), Inches(0.4),
         [(f"{author}, {c['licence']}", 16, False, INK2)])
    y += Inches(0.62)
text(s, MARGIN, y + Inches(0.35), W - 2 * MARGIN, Inches(0.9),
     [[("Satellite imagery and 80 × 80 chips", 16, True, INK),
       (": Ships in Satellite Imagery (ShipsNet), Planet Labs / Kaggle, CC BY-SA 4.0.",
        16, False, INK2)]], line=1.3)

prs.save(OUT)
print(f"wrote {OUT.relative_to(HERE.parent)}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")

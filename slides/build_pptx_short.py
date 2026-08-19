"""Build the short (20-slide) ICMIC 2026 PowerPoint deck for CONFROUTE.

Same story, same figures, same palette as build_pptx.py, compressed to 20 slides.
Self-contained on purpose: the full 31-slide deck keeps building from
build_pptx.py, untouched.

What is merged relative to the full deck:
  * the four scene-setting photographs become two
  * the question and the assumption share one statement slide
  * the nominal-only chip slide is dropped, fig_chips already carries both rows
  * "a gate fails exactly when you need it" becomes the title of the error slide
  * the CONFROUTE section break folds into the architecture slide
  * the ablation and the action mix share one slide
  * image credits move to a band on the closing slide
The backup slides (full results table, confidence traces) live only in the full
deck. Keep it open in a second window for Q&A.

    python slides/make_figures.py && python slides/build_pptx_short.py
"""

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from PIL import Image
from lxml import etree

HERE = Path(__file__).resolve().parent
FIG = HERE / "figures"
IMG = HERE / "images"
ARCH = HERE.parent / "paper" / "figures" / "architecture_full.png"
OUT = HERE / "CONFROUTE_ICMIC2026_short.pptx"

W, H = Inches(13.333), Inches(7.5)

# palette, identical to the figures
NAVY = RGBColor(0x10, 0x46, 0x80)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
AQUA = RGBColor(0x1B, 0xAF, 0x7A)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
RULE = RGBColor(0xE1, 0xE0, 0xD9)
WASH = RGBColor(0xF4, 0xF7, 0xFB)
WASH_O = RGBColor(0xFD, 0xF0, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"
MARGIN = Inches(0.72)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------------------- #
# primitives
# --------------------------------------------------------------------------- #
def slide(notes=None):
    s = prs.slides.add_slide(BLANK)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def _alpha(color_elm, alpha):
    """python-pptx has no transparency API; splice an <a:alpha> into the fill."""
    srgb = color_elm.find(qn("a:srgbClr"))
    if srgb is None:
        return
    a = etree.SubElement(srgb, qn("a:alpha"))
    a.set("val", str(int(alpha * 100000)))


def rect(s, x, y, w, h, color, alpha=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if alpha < 1.0:
        _alpha(shp.fill._xPr.find(qn("a:solidFill")), alpha)
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line=1.15, space_after=0):
    """runs: list of (string, size_pt, bold, color) or a list of such lists (paragraphs)."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = runs if isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(space_after)
        for content, size, bold, color in para:
            r = p.add_run()
            r.text = content
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = FONT
    return box


def pic_cover(s, path, x, y, w, h):
    """Fill the box, preserving aspect ratio by centre-cropping (CSS object-fit: cover)."""
    iw, ih = Image.open(path).size
    box_ar, img_ar = w / h, iw / ih
    pic = s.shapes.add_picture(str(path), x, y, w, h)
    if img_ar > box_ar:                       # image too wide -> crop sides
        keep = box_ar / img_ar
        pic.crop_left = pic.crop_right = (1 - keep) / 2
    else:                                     # image too tall -> crop top/bottom
        keep = img_ar / box_ar
        pic.crop_top = pic.crop_bottom = (1 - keep) / 2
    return pic


def pic_fit(s, path, x, y, w, h):
    """Contain the whole image inside the box, centred."""
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    nw, nh = int(iw * scale), int(ih * scale)
    return s.shapes.add_picture(str(path), x + int((w - nw) / 2), y + int((h - nh) / 2), nw, nh)


def footer(s, n):
    text(s, MARGIN, H - Inches(0.52), Inches(8), Inches(0.3),
         [("CONFROUTE  |  ICMIC 2026, NTUST Taipei", 10, False, MUTED)])
    text(s, W - MARGIN - Inches(1), H - Inches(0.52), Inches(1), Inches(0.3),
         [(str(n), 10, False, MUTED)], align=PP_ALIGN.RIGHT)


# --------------------------------------------------------------------------- #
# slide templates
# --------------------------------------------------------------------------- #
def photo_slide(headline, sub=None, image=None, notes=None, accent=None, side="bottom"):
    """Full-bleed photograph with a scrim and one short headline."""
    s = slide(notes)
    pic_cover(s, IMG / image, 0, 0, W, H)
    if side == "bottom":
        rect(s, 0, Inches(3.9), W, H - Inches(3.9), RGBColor(0x05, 0x14, 0x24), 0.72)
        y = Inches(4.55)
    else:
        rect(s, 0, 0, W, H, RGBColor(0x05, 0x14, 0x24), 0.55)
        y = Inches(2.7)
    text(s, MARGIN, y, W - 2 * MARGIN, Inches(1.6),
         [(headline, 40, True, WHITE)], line=1.1)
    if sub:
        gap = Inches(1.35) if side == "bottom" else Inches(1.75)
        text(s, MARGIN, y + gap, W - 2 * MARGIN, Inches(1.0),
             [(sub, 20, False, RGBColor(0xD8, 0xE4, 0xF2))], line=1.25)
    if accent:
        rect(s, MARGIN, y - Inches(0.34), Inches(1.1), Pt(5), accent or BLUE)
    return s


def statement_slide(headline, sub=None, notes=None, color=NAVY):
    """White slide, one sentence, nothing else."""
    s = slide(notes)
    text(s, Inches(1.4), Inches(2.5), W - Inches(2.8), Inches(1.9),
         [(headline, 44, True, color)], align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.BOTTOM, line=1.12)
    if sub:
        text(s, Inches(2.1), Inches(4.6), W - Inches(4.2), Inches(1.4),
             [(sub, 20, False, INK2)], align=PP_ALIGN.CENTER, line=1.35)
    return s


def figure_slide(title, image, n, sub=None, note=None, notes=None):
    """A chart or diagram with the takeaway as the title."""
    s = slide(notes)
    text(s, MARGIN, Inches(0.5), W - 2 * MARGIN, Inches(0.8),
         [(title, 30, True, NAVY)], line=1.1)
    top = Inches(1.32)
    if sub:
        text(s, MARGIN, Inches(1.22), W - 2 * MARGIN, Inches(0.45),
             [(sub, 16, False, INK2)])
        top = Inches(1.78)
    bottom = H - Inches(1.35) if note else H - Inches(0.75)
    pic_fit(s, image, MARGIN, top, W - 2 * MARGIN, bottom - top)
    if note:
        rect(s, MARGIN, H - Inches(1.24), W - 2 * MARGIN, Inches(0.72), WASH_O)
        text(s, MARGIN + Inches(0.28), H - Inches(1.13), W - 2 * MARGIN - Inches(0.56),
             Inches(0.5), note, anchor=MSO_ANCHOR.MIDDLE, line=1.2)
    footer(s, n)
    return s


def title_only(title, n, sub=None, notes=None):
    s = slide(notes)
    text(s, MARGIN, Inches(0.5), W - 2 * MARGIN, Inches(0.8),
         [(title, 30, True, NAVY)], line=1.1)
    if sub:
        text(s, MARGIN, Inches(1.22), W - 2 * MARGIN, Inches(0.45),
             [(sub, 16, False, INK2)])
    footer(s, n)
    return s


# --------------------------------------------------------------------------- #
# the deck, 20 slides
# --------------------------------------------------------------------------- #
n = 0


def num():
    global n
    n += 1
    return n


# ---- 1. title -------------------------------------------------------------
s = slide("Name, affiliation, one line: this is a paper about when an edge AI "
          "system on a ship should refuse to make a decision. 20 s, do not linger.")
pic_cover(s, IMG / "hero.jpg", 0, 0, W, H)
rect(s, 0, 0, W, H, RGBColor(0x04, 0x11, 0x20), 0.62)
text(s, MARGIN, Inches(1.05), W - 2 * MARGIN, Inches(0.4),
     [("ICMIC 2026  ·  NTUST, Taipei  ·  19–21 August 2026", 14, False,
       RGBColor(0x9E, 0xC5, 0xF4))])
text(s, MARGIN, Inches(1.95), W - 2 * MARGIN, Inches(1.1),
     [("CONFROUTE", 66, True, WHITE)], line=1.0)
text(s, MARGIN, Inches(3.0), W - Inches(4.0), Inches(1.3),
     [("Confidence-Aware Routing for Resilient Maritime Edge AI", 28, False,
       RGBColor(0xCD, 0xE2, 0xFB))], line=1.2)
rect(s, MARGIN, Inches(4.45), Inches(1.6), Pt(5), BLUE)
text(s, MARGIN, Inches(4.95), W - 2 * MARGIN, Inches(0.9),
     [[("Anthony Uchenna Eneh", 19, True, WHITE),
       ("   ·   Love Allen Chijioke Ahakonye   ·   Jae Min Lee   ·   "
        "Dong-Seong Kim", 19, False, RGBColor(0xD8, 0xE4, 0xF2))]], line=1.3)
text(s, MARGIN, Inches(5.65), W - 2 * MARGIN, Inches(1.0),
     [("Kumoh National Institute of Technology, Gumi, South Korea   ·   "
       "NSLab Co. Ltd.", 14, False, RGBColor(0x9E, 0xC5, 0xF4))], line=1.35)

# ---- 2. the environment, in one slide ------------------------------------
photo_slide("Ships now run AI on board.",
            "Perception, monitoring, navigation support, on a satellite link that is "
            "intermittent by nature.",
            image="sensor.jpg", accent=BLUE,
            notes="Two facts, fast: the compute is already on the vessel, and the link "
                  "that would let it ask for help is not guaranteed. Offloading is not "
                  "an escape hatch. ~30 s.")
footer(prs.slides[-1], num())

# ---- 3. the input drifts, and errors outrank latency ---------------------
photo_slide("And the sea changes what the camera sees.",
            "Sea state, haze, illumination, occlusion: the input drifts away from "
            "training, and at sea a wrong answer costs more than a slow one.",
            image="fog.jpg", accent=ORANGE, side="full",
            notes="Point at the haze. Then land the inversion: latency-optimal is not "
                  "safety-optimal, which is what separates this from ordinary MEC work. "
                  "Say it plainly and pause. ~35 s.")
footer(prs.slides[-1], num())

# ---- 4. the question, and the assumption ---------------------------------
statement_slide("Act now, ask for help, or refuse to answer?",
                "Every task, on every node, under uncertainty. Confidence-threshold "
                "offloading assumes the model knows when it is unsure. We tested that.",
                notes="The question the paper answers, and the assumption the field is "
                      "built on. Do not defend the assumption, you are about to break "
                      "it. ~30 s.")
footer(prs.slides[-1], num())

# ---- 5. the measurement, on real chips -----------------------------------
figure_slide("It finds every ship, until you add haze.",
             FIG / "fig_chips.png", num(),
             sub="Three held-out ShipsNet chips, 80 × 80 PlanetScope imagery, "
                 "before and after the adverse transform.",
             note=[("Every chip still contains a ship. The model is now ", 15, False, INK),
                   ("wrong on all three, and far more confident", 15, True, ORANGE),
                   (".", 15, False, INK)],
             notes="THE slide. Top row: all three correct, and confidence is only 0.51 "
                   "to 0.67, honestly unsure. Bottom row: same three vessels, NO SHIP at "
                   "0.96, 0.88, 0.91. \"It did not become uncertain. It became more "
                   "certain, and wrong.\" Pause. ~1:10.")

# ---- 6. the aggregate -----------------------------------------------------
figure_slide("It is not three unlucky chips.",
             FIG / "fig_calibration.png", num(),
             sub="500 held-out images, nominal versus adverse.",
             note=[("Accuracy falls 13.8 points. Mean confidence ", 15, False, INK),
                   ("rises", 15, True, ORANGE),
                   (" 3.2 points.", 15, False, INK)],
             notes="The population-level version of the previous slide. The lines cross. "
                   "That crossing is the paper. Those three chips are 3 of 53 that behave "
                   "this way. ~50 s.")

# ---- 7. the consequence ---------------------------------------------------
figure_slide("A confidence gate fails exactly when you need it.",
             FIG / "fig_error_confidence.png", num(),
             note=[("Under adverse conditions ", 15, False, INK),
                   ("90% of errors sit above the 0.65 threshold", 15, True, ORANGE),
                   (", invisible to it. Under nominal conditions, 39%.", 15, False, INK)],
             notes="Errors only, plotted against the fallback threshold. The gate is "
                   "blind precisely in the regime that threatens the vessel. This is the "
                   "hinge into the method. ~45 s.")

# ---- 8. the method --------------------------------------------------------
figure_slide("CONFROUTE: make refusal a routing action.",
             ARCH, num(),
             sub="Post-scoring, pre-actuation: consumes uncertainty without touching "
                 "the model, logs every decision.",
             notes="Trace the three branches with the pointer. Post-scoring, "
                   "pre-actuation, so it drops into an existing stack. ~45 s.")

# ---- 9. three actions -----------------------------------------------------
s = slide("Local when confident and the queue allows. Offload when the link is good and "
          "the queue is loaded. Fallback when uncertain AND cut off. Land the guardrail "
          "line: the threshold controller is inspectable and certifiable. ~40 s.")
text(s, MARGIN, Inches(0.5), W - 2 * MARGIN, Inches(0.8),
     [("Three actions, one of them new.", 30, True, NAVY)], line=1.1)
cards = [
    ("LOCAL", BLUE, "Run the model here.", "Confident enough to act alone."),
    ("OFFLOAD", AQUA, "Ask a peer or relay.", "Link is usable, queue is loaded."),
    ("FALLBACK", ORANGE, "Refuse, and act safely.", "Uncertain and cut off."),
]
cw, gap = Inches(3.75), Inches(0.42)
x0 = (W - (3 * cw + 2 * gap)) / 2
for i, (name, colour, line1, line2) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rect(s, x, Inches(1.85), cw, Inches(3.5), WASH if i < 2 else WASH_O)
    rect(s, x, Inches(1.85), cw, Pt(7), colour)
    text(s, x + Inches(0.4), Inches(2.35), cw - Inches(0.8), Inches(0.6),
         [(name, 26, True, colour)])
    text(s, x + Inches(0.4), Inches(3.15), cw - Inches(0.8), Inches(1.0),
         [(line1, 19, True, INK)], line=1.2)
    text(s, x + Inches(0.4), Inches(3.95), cw - Inches(0.8), Inches(1.2),
         [(line2, 16, False, INK2)], line=1.25)
text(s, MARGIN, Inches(5.75), W - 2 * MARGIN, Inches(0.6),
     [[("Scores decide, ", 17, False, INK2),
       ("guardrails override", 17, True, NAVY),
       (": no offload without a peer, no local execution on a safety-critical task "
        "below the confidence floor.", 17, False, INK2)]], line=1.3)
footer(s, num())

# ---- 10. fallback is concrete --------------------------------------------
s = title_only("Fallback is an action, not an error code.", num(),
               sub="What the conservative branch means, per maritime application.",
               notes="Kills the objection that fallback means crash or do nothing. Read "
                     "the collision-avoidance row aloud, then move on. ~35 s.")
rows = [
    ("Collision avoidance", "Reduce speed, sound horn, hold heading, alert the operator"),
    ("Maritime surveillance", "Flag the chip for review, defer the alert, request another pass"),
    ("Engine anomaly", "Report “inspection needed”, fall back to a limit-based alert"),
    ("Waypoint decision", "Hold position, revert to last safe heading, request human review"),
]
y = Inches(2.05)
for i, (app, act) in enumerate(rows):
    if i % 2 == 0:
        rect(s, MARGIN, y - Inches(0.12), W - 2 * MARGIN, Inches(0.95), WASH)
    text(s, MARGIN + Inches(0.3), y, Inches(3.6), Inches(0.7),
         [(app, 19, True, NAVY)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, MARGIN + Inches(4.1), y, W - 2 * MARGIN - Inches(4.5), Inches(0.7),
         [(act, 18, False, INK)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.06)
text(s, MARGIN, Inches(6.42), W - 2 * MARGIN, Inches(0.5),
     [("Conservative, auditable, and aligned with COLREGS/IMO practice: "
       "the ship does something safe and says why.", 15, False, INK2)])

# ---- 11. when fallback fires ---------------------------------------------
figure_slide("Refuse only when both signals fail.",
             FIG / "fig_boundary.png", num(),
             sub="Confidence below 0.65 and bandwidth below 1.8 Mbps.",
             note=[("If the link is healthy, the router would rather ", 15, False, INK),
                   ("ask a better-resourced peer than refuse", 15, True, NAVY),
                   (". That keeps fallback rare.", 15, False, INK)],
             notes="Requiring both conditions is what stops the router hiding behind "
                   "fallback. ~30 s.")

# ---- 12. setup ------------------------------------------------------------
s = slide("ShipsNet, ResNet-18, three link regimes, four baselines. Say the unsafe-rate "
          "definition out loud: routing-layer risk, not physical incidents. ~35 s.")
text(s, MARGIN, Inches(0.5), W - 2 * MARGIN, Inches(0.8),
     [("37,500 routing decisions.", 30, True, NAVY)], line=1.1)
text(s, MARGIN, Inches(1.22), W - 2 * MARGIN, Inches(0.45),
     [("2,500 tasks × 5 policies × 3 link regimes.", 16, False, INK2)])
stats = [
    ("4,000", "PlanetScope chips\n80 × 80, ShipsNet"),
    ("ResNet-18", "ImageNet-initialised\n84.8% validation"),
    ("3", "link regimes\nstable / intermittent / degraded"),
    ("4", "baselines\nlocal, offload, confidence, load"),
]
cw, gap = Inches(2.85), Inches(0.32)
x0 = (W - (4 * cw + 3 * gap)) / 2
for i, (big, small) in enumerate(stats):
    x = x0 + i * (cw + gap)
    rect(s, x, Inches(2.1), cw, Inches(2.6), WASH)
    text(s, x, Inches(2.5), cw, Inches(0.9), [(big, 40, True, BLUE)],
         align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.25), Inches(3.5), cw - Inches(0.5), Inches(1.1),
         [(small, 15, False, INK2)], align=PP_ALIGN.CENTER, line=1.3)
text(s, MARGIN, Inches(5.25), W - 2 * MARGIN, Inches(1.0),
     [[("Unsafe rate", 17, True, NAVY),
       (" = routing-layer risk: safety-sensitive local errors, low-confidence local "
        "decisions, failed safety-sensitive offloads.\nIt is not a count of physical "
        "incidents.", 17, False, INK2)]], line=1.35)
footer(s, num())

# ---- 13-15. results -------------------------------------------------------
figure_slide("Every baseline converges at about 7%.",
             FIG / "fig_unsafe_by_scenario.png", num(),
             note=[("When the link goes, so does their escape hatch. ", 15, False, INK),
                   ("CONFROUTE holds at 4.8%", 15, True, BLUE),
                   (".", 15, False, INK)],
             notes="The convergence is the story, not our line. Stable links flatter "
                   "everything that offloads; degraded links expose the dependency. Then "
                   "point at 4.8%. ~55 s.")

figure_slide("Degraded links: safer and faster.",
             FIG / "fig_tradeoff.png", num(),
             note=[("Unsafe 6.8% → 4.8% and latency 98.2 → 58.5 ms against "
                    "always-offload. ", 15, False, INK),
                   ("A failed offload costs both.", 15, True, BLUE)],
             notes="Normally safety costs latency. Not here, because attempting an "
                   "offload that fails wastes time AND leaves the task unhandled. ~50 s.")

s = title_only("Where we do not win: intermittent links.", num(),
               sub="The regime in which offloading still mostly works.",
               notes="Deliberate honesty, and it buys credibility. Load-only is safer "
                     "here because offloading still works and removes local-error risk. "
                     "Say why, then land the closing line. ~40 s.")
tbl_rows = [
    ("Always local", "49.4", "5.8%", "0.0%", False),
    ("Always offload", "71.6", "4.3%", "0.0%", False),
    ("Confidence threshold", "59.6", "4.9%", "0.0%", False),
    ("Load only", "67.0", "3.8%", "0.0%", False),
    ("CONFROUTE (ours)", "61.6", "4.4%", "2.6%", True),
]
hx = [MARGIN + Inches(0.3), MARGIN + Inches(4.6), MARGIN + Inches(6.7), MARGIN + Inches(8.9)]
for label, x in zip(["Policy", "Latency (ms)", "Unsafe rate", "Fallback"], hx):
    text(s, x, Inches(2.0), Inches(4.2), Inches(0.4), [(label, 16, True, NAVY)])
rect(s, MARGIN + Inches(0.3), Inches(2.45), W - 2 * MARGIN - Inches(0.6), Pt(1.5), RULE)
y = Inches(2.62)
for label, lat, uns, fb, ours in tbl_rows:
    if ours:
        rect(s, MARGIN, y - Inches(0.1), W - 2 * MARGIN, Inches(0.62), WASH)
    c = BLUE if ours else INK
    for val, x in zip([label, lat, uns, fb], hx):
        text(s, x, y, Inches(4.2), Inches(0.5), [(val, 19, ours, c)])
    y += Inches(0.66)
rect(s, MARGIN, Inches(6.05), W - 2 * MARGIN, Inches(0.92), WASH_O)
text(s, MARGIN + Inches(0.3), Inches(6.2), W - 2 * MARGIN - Inches(0.6), Inches(0.7),
     [[("Load-only routing is safer here", 16, True, INK),
       (" (3.8% vs 4.4%), at 5.4 ms more latency. Our case is the degraded regime, "
        "which is also the regime that actually threatens the vessel.", 16, False, INK)]],
     line=1.3)

# ---- 16. ablation and action mix, together --------------------------------
s = title_only("Both parts are load-bearing, and the router adapts.", num(),
               sub="Unsafe outcome rate under degraded links, and the action mix by regime.",
               notes="Removing fallback hurts more than removing link awareness: having "
                     "somewhere to put an undecidable task is the bigger effect. Then the "
                     "chart: offload collapses 89% to 11% as the link degrades, and "
                     "fallback stays a bounded minority. ~45 s.")
abl = [("Full CONFROUTE", "4.8%", "", True),
       ("No fallback action", "6.2%", "p = 0.012", False),
       ("No link-state awareness", "5.9%", "p = 0.037", False)]
y = Inches(2.35)
for label, val, sig, ours in abl:
    if ours:
        rect(s, MARGIN, y - Inches(0.16), Inches(4.95), Inches(0.9), WASH)
    text(s, MARGIN + Inches(0.25), y, Inches(3.1), Inches(0.5),
         [(label, 18, ours, BLUE if ours else INK)])
    text(s, MARGIN + Inches(3.5), y - Inches(0.06), Inches(1.4), Inches(0.6),
         [(val, 24, True, BLUE if ours else INK)])
    text(s, MARGIN + Inches(0.25), y + Inches(0.42), Inches(3.1), Inches(0.4),
         [(sig, 14, False, INK2)])
    y += Inches(1.12)
text(s, MARGIN + Inches(0.25), Inches(5.95), Inches(4.7), Inches(0.9),
     [("Fallback carries more of the benefit than link awareness does.",
       15, False, INK2)], line=1.3)
pic_fit(s, FIG / "fig_actions.png", MARGIN + Inches(5.35), Inches(2.1),
        W - 2 * MARGIN - Inches(5.35), Inches(3.55))
text(s, MARGIN + Inches(5.35), Inches(5.95), W - 2 * MARGIN - Inches(5.35), Inches(0.9),
     [[("Offload collapses 89% → 11%. Fallback stays a ", 15, False, INK2),
       ("bounded minority", 15, True, NAVY),
       (": the vessel keeps operating.", 15, False, INK2)]], line=1.3)

# ---- 17. the point --------------------------------------------------------
statement_slide("The routing layer was\nnever the bottleneck.",
                "We improved it, and unsafe outcomes still occur at high confidence.",
                notes="Slow down. Let it sit before speaking. This is what you want the "
                      "room to remember. ~30 s.")
footer(prs.slides[-1], num())

# ---- 18. limitations ------------------------------------------------------
s = title_only("What we are not claiming.", num(),
               notes="State the limitations without apologising, then the forward path "
                     "with calibration first. ~40 s.")
cols = [
    # one line each: wrapped bullets lose their hanging indent
    ("Limitations", ORANGE, [
        "Degradation is synthetic.",
        "One task family: satellite chips.",
        "Outcomes from a cost model.",
        "Thresholds fixed and trace-tuned.",
    ]),
    ("Where this goes", NAVY, [
        "Calibration first: conformal, OOD.",
        "Adaptive thresholds via bandits.",
        "Peer-aware forwarding.",
        "Real traces, hardware-in-the-loop.",
    ]),
]
cw = Inches(5.7)
for i, (head, colour, items) in enumerate(cols):
    x = MARGIN + i * (cw + Inches(0.7))
    text(s, x, Inches(1.75), cw, Inches(0.5), [(head, 22, True, colour)])
    rect(s, x, Inches(2.28), Inches(1.0), Pt(4), colour)
    y = Inches(2.65)
    for it in items:
        text(s, x, y, cw, Inches(1.0), [("•  " + it, 17, False, INK)], line=1.3)
        y += Inches(1.02)

# ---- 19. takeaways --------------------------------------------------------
s = title_only("Three things to take away.", num(),
               notes="Read the three headers only, then stop and invite questions. Do not "
                     "add a new thought here. ~50 s.")
takeaways = [
    ("Fallback deserves to be a routing action.",
     "Under degraded links it cuts unsafe outcomes from 6.8% to 4.8% and saves 40 ms."),
    ("Confidence thresholds break exactly when needed.",
     "Accuracy falls 13.8 points while confidence rises, and 90% of errors clear the gate."),
    ("Calibration is the reliability bottleneck.",
     "Ahead of routing or scheduling. Safe autonomy at sea depends on knowing "
     "what the model does not know."),
]
y = Inches(1.9)
for i, (head, body) in enumerate(takeaways):
    text(s, MARGIN, y, Inches(0.7), Inches(0.7), [(str(i + 1), 34, True, BLUE)])
    text(s, MARGIN + Inches(0.85), y + Inches(0.03), W - 2 * MARGIN - Inches(0.9),
         Inches(0.5), [(head, 23, True, INK)], line=1.15)
    text(s, MARGIN + Inches(0.85), y + Inches(0.62), W - 2 * MARGIN - Inches(0.9),
         Inches(0.8), [(body, 17, False, INK2)], line=1.3)
    y += Inches(1.55)

# ---- 20. thank you, with the image credits band --------------------------
# credits are generated from the fetch manifest so attribution cannot drift
def tidy_author(raw):
    """Commons artist fields for derivative works carry filename noise."""
    a = " ".join((raw or "see source").split())
    a = re.sub(r"^File:.*?\.(jpg|jpeg|png|tif+)\s*:\s*", "", a, flags=re.I)
    a = re.sub(r"\s*derivative work\s*:\s*", " / ", a, flags=re.I)
    a = re.sub(r"^(User|user):\s*", "", a)
    return a[:60].strip()


# only the photographs this shorter deck actually shows
USED_SLOTS = ("hero", "sensor", "fog")
credits = [c for c in json.loads((IMG / "CREDITS.json").read_text(encoding="utf-8"))
           if c["slot"] in USED_SLOTS]
credit_line = ("Photographs, Wikimedia Commons: "
               + "; ".join(f"{tidy_author(c['artist'])} ({c['licence']})"
                           for c in credits)
               + ".  Satellite imagery and 80 × 80 chips: Ships in Satellite "
                 "Imagery (ShipsNet), Planet Labs / Kaggle, CC BY-SA 4.0.")

s = slide("Thank you, and take questions. The full 31-slide deck has the backup "
          "material: full results table, confidence traces, per-image credits.")
pic_cover(s, IMG / "hero.jpg", 0, 0, W, H)
rect(s, 0, 0, W, H, RGBColor(0x04, 0x11, 0x20), 0.66)
text(s, MARGIN, Inches(2.35), W - 2 * MARGIN, Inches(1.2),
     [("Thank you. Questions?", 48, True, WHITE)], line=1.1)
rect(s, MARGIN, Inches(3.7), Inches(1.6), Pt(5), BLUE)
text(s, MARGIN, Inches(4.15), W - 2 * MARGIN, Inches(1.6),
     [[("Anthony Uchenna Eneh\n", 20, True, WHITE),
       ("anthony@kumoh.ac.kr\nKumoh National Institute of Technology, Gumi, South Korea",
        18, False, RGBColor(0xD8, 0xE4, 0xF2))]], line=1.4)
rect(s, 0, Inches(6.15), W, H - Inches(6.15), RGBColor(0x03, 0x0C, 0x18), 0.45)
text(s, MARGIN, Inches(6.42), W - 2 * MARGIN, Inches(0.7),
     [(credit_line, 10.5, False, RGBColor(0xC3, 0xD4, 0xE8))], line=1.35)
num()

prs.save(OUT)
print(f"wrote {OUT.relative_to(HERE.parent)}  "
      f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
